import collections 
import collections.abc
from pptx import Presentation
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Pt,Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import XySeriesData, XyChartData
import chart
import os
# # Create presentation and slide
# prs = Presentation()
# slide = prs.slides.add_slide(prs.slide_layouts[5])

# chart_data = ChartData()
# chart_data.categories = ['Fruits', 'Vegetables', 'Meat']
# chart_data.add_series('Dummy Data', (35, 25, 40))

# # Add doughnut chart
# chart = slide.shapes.add_chart(
#     XL_CHART_TYPE.DOUGHNUT, Inches(2),Inches(2),Inches(2),Inches(2) , chart_data
# ).chart

# # Add dummy data to chart
# chart.replace_data(chart_data)

# # Save the presentation
# prs.save('doughnut_chart_dummy_data.pptx')

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

# Create presentation and slide
prs = Presentation()

blank_slide_layout = prs.slide_layouts[6]
slide2 = prs.slides.add_slide(blank_slide_layout)
chart.add_chart(slide=slide2, chartName="COLUMN_CLUSTERED")

blank_slide_layout = prs.slide_layouts[6]
slide3 = prs.slides.add_slide(blank_slide_layout)
chart.add_chart(slide=slide3, chartName="COLUMN_STACKED")

blank_slide_layout = prs.slide_layouts[6]
slide4 = prs.slides.add_slide(blank_slide_layout)
chart.add_chart(slide=slide4, chartName="COLUMN_STACKED_100")

blank_slide_layout = prs.slide_layouts[6]
slide5 = prs.slides.add_slide(blank_slide_layout)
chart.add_chart(slide=slide5, chartName="LINE")

blank_slide_layout = prs.slide_layouts[6]
slide6 = prs.slides.add_slide(blank_slide_layout)
chart.add_chart(slide=slide6, chartName="LINE_STACKED")

blank_slide_layout = prs.slide_layouts[6]
slide7 = prs.slides.add_slide(blank_slide_layout)
chart.add_chart(slide=slide7, chartName="LINE_MARKERS_STACKED_100")

blank_slide_layout = prs.slide_layouts[6]
slide8 = prs.slides.add_slide(blank_slide_layout)
chart.add_chart(slide=slide8, chartName="LINE_MARKERS_STACKED")

blank_slide_layout = prs.slide_layouts[6]
slide9 = prs.slides.add_slide(blank_slide_layout)
chart.add_chart(slide=slide9, chartName="LINE_STACKED_100")

blank_slide_layout = prs.slide_layouts[6]
slide9 = prs.slides.add_slide(blank_slide_layout)
chart.add_chart(slide=slide9, chartName="LINE_MARKERS")

chartName = ["PIE","PIE_EXPLODED","DOUGHNUT","DOUGHNUT_EXPLODED","BAR_CLUSTERED","BAR_STACKED","BAR_STACKED_100","AREA","AREA_STACKED","AREA_STACKED_100"]

for i in chartName:
    blank_slide_layout = prs.slide_layouts[6]
    slide9 = prs.slides.add_slide(blank_slide_layout)
    chart.add_chart(slide=slide9, chartName=i)


prs.save('test3.pptx')
os.startfile('test3.pptx')