import collections 
import collections.abc
from pptx import Presentation
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Pt,Inches
# from pptx.enum.dml import MSO_SHADOW_STYLE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.chart.data import XySeriesData, XyChartData
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.chart import XL_TICK_MARK
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_MARKER_STYLE
from pptx.enum.text import PP_ALIGN
import os
from pptx.enum.chart import XL_LABEL_POSITION

prs = Presentation()

prs.slide_width = Inches(10.833)
prs.slide_height = Inches(7.5)

blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

def create_chart(widget:dict):
    horizontal_position , vertical_position , width , height = Inches(widget["horizontal_position"]),Inches(widget["vertical_position"]),Inches(widget["width"]),Inches(widget["height"])
    
    chart = widget["widget"]["value"]

    if chart["chart_type"] in ["line","bar","column"]:
        chartOption = chart["options"]
        if chart["chart_type"] == "line":
            chart_data = ChartData()
            categories = chart["chart_data"]["catagories_name"]
            chart_data.categories = categories
            for series in chart["chart_data"]["series"]:
                chart_data.add_series(series["series_name"], series["data"])
            
            chart_option = ''
            if chartOption.get("line_type"):
                if chartOption["line_type"] == "LINE_MARKERS":
                    charts = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, horizontal_position, vertical_position, width, height, chart_data)
                    chart_option = charts.chart
                elif chartOption["line_type"] == "LINE_MARKERS_STACKED":
                    charts = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS_STACKED, horizontal_position, vertical_position, width, height, chart_data)
                    chart_option = charts.chart
                elif chartOption["line_type"] == "LINE_MARKERS_STACKED_100":
                    charts = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS_STACKED_100, horizontal_position, vertical_position, width, height, chart_data)
                    chart_option = charts.chart
                elif chartOption["line_type"] == "LINE_STACKED":
                    charts = slide.shapes.add_chart(XL_CHART_TYPE.LINE_STACKED, horizontal_position, vertical_position, width, height, chart_data)
                    chart_option = charts.chart
                elif chartOption["line_type"] == "LINE_STACKED_100":
                    charts = slide.shapes.add_chart(XL_CHART_TYPE.LINE_STACKED_100, horizontal_position, vertical_position, width, height, chart_data)
                    chart_option = charts.chart
                else:
                    charts = slide.shapes.add_chart(XL_CHART_TYPE.LINE, horizontal_position, vertical_position, width, height, chart_data)
                    chart_option = charts.chart
            else:
                charts = slide.shapes.add_chart(XL_CHART_TYPE.LINE, horizontal_position, vertical_position, width, height, chart_data)
                chart_option = charts.chart

            chart_option.has_title = False
            if chartOption.get("title"):
                chart_option.has_title = True
            
            if chartOption.get("legend"):
                legend = chartOption["legend"]
                chart_option.has_legend = True
                if legend["legend_potision"] == "left":
                    chart_option.legend.position = XL_LEGEND_POSITION.LEFT
                elif legend["legend_potision"] == "right":
                    chart_option.legend.position = XL_LEGEND_POSITION.RIGHT
                elif legend["legend_potision"] == "TOP":
                    chart_option.legend.position = XL_LEGEND_POSITION.TOP
                elif legend["legend_potision"] == "bottom":
                    chart_option.legend.position = XL_LEGEND_POSITION.BOTTOM
                else:
                    chart_option.legend.position = XL_LEGEND_POSITION.LEFT

                if legend.get("font_style_setting"):
                    set_chart_text(chart_option.legend, legend["font_style_setting"])

            if chartOption.get("value_axis"):
                value_axis = chartOption["value_axis"]
                axis = chart_option.value_axis

                if value_axis.get("minimum_scale"):
                    axis.minimum_scale = value_axis["minimum_scale"]
                if value_axis.get("maximum_scale"):
                    axis.maximum_scale = value_axis["maximum_scale"]
                
                if value_axis.get("format_number"):
                    if value_axis.get("format_number") == "number":
                        axis.tick_labels.number_format = '#,###'
                    elif value_axis.get("format_number") == "presentage":
                        axis.tick_labels.number_format = '0.0%'
                    elif value_axis.get("format_number") == "date":
                        axis.tick_labels.number_format = 'dd/mm/yyyy'

                axis.has_major_gridlines = True
                axis.has_minor_gridlines = True

                if value_axis.get("major_gridlines"):
                    axis.has_major_gridlines = True
                else:
                    axis.has_major_gridlines = False

                if value_axis.get("minor_gridlines"):
                    axis.has_minor_gridlines = True
                else:
                    axis.has_minor_gridlines = False

            if chartOption.get("category_axis"):
                category_axis = chartOption["category_axis"]

                cate_asix = chart_option.category_axis
                cate_asix.has_major_gridlines = True
                cate_asix.has_minor_gridlines = True

                if category_axis.get("major_gridlines"):
                    cate_asix.has_major_gridlines = True
                else:
                    cate_asix.has_major_gridlines = False

                if category_axis.get("minor_gridlines"):
                    cate_asix.has_minor_gridlines = True
                else:
                    cate_asix.has_minor_gridlines = False

            if chartOption.get("series"):
                series = chartOption["series"]
                if series.get("color"):
                    plot = chart_option.plots[0]
                    serie = plot.series

                    for i in range(len(serie)):
                        line = serie[i].format.line
                        serie[i].marker.format.fill.solid()
                        serie[i].marker.format.fill.fore_color.rgb = RGBColor.from_string(series.get("color")[i])
                        serie[i].marker.style = XL_MARKER_STYLE.CIRCLE
                        line.color.rgb = RGBColor.from_string(series.get("color")[i])


def set_chart_text(textframe, text_config):
    if text_config.get("font_size"):
            textframe.font.size = Pt(text_config["font_size"])
    else:
        textframe.font.size = Pt(14)
    if text_config.get("bold") == True:
        textframe.font.bold = True

    if text_config.get("underline") == True:
        textframe.font.underline = True

    if text_config.get("italic") == True:
        textframe.font.italic = True

    if text_config.get("font_color_hex"):
        textframe.font.color.rgb = RGBColor.from_string(text_config["font_color_hex"])

def create_content(content:dict):
    horizontal_position , vertical_position , width , height = Inches(content["horizontal_position"]),Inches(content["vertical_position"]),Inches(content["width"]),Inches(content["height"])

    txBox = slide.shapes.add_textbox(horizontal_position, vertical_position, width, height)

    text_frame = txBox.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]

    for text in content["content"]:
        option = text["option"]
        run = p.add_run()

        if option.get("font_size"):
            run.font.size = Pt(option["font_size"])
        else:
            run.font.size = Pt(14)

        if option.get("bold") == True:
            run.font.bold = True

        if option.get("underline") == True:
            run.font.underline = True

        if option.get("italic") == True:
            run.font.italic = True

        if option.get("font_color_hex"):
            run.font.color.rgb = RGBColor.from_string(option["font_color_hex"])
        
        if content.get("alignment"):
            if content["alignment"] == "left":
                p.alignment = PP_ALIGN.LEFT
            elif content["alignment"] == "right":
                p.alignment = PP_ALIGN.RIGHT
            elif content["alignment"] == "center":
                p.alignment = PP_ALIGN.CENTER
            else:
                p.alignment = PP_ALIGN.JUSTIFY
        
        run.text = text["text"]

content1= {
  "height":0.93,
  "width": 7.8,
  "horizontal_position":0.06,
  "vertical_position": 0.12,
  "content":[
    {
      "text": "Ketenagakerjaan",
      "alignment": "left",
      "option":{
        "font_size" : 44,
        "bold": True
      }
    }
  ]
}
content2= {
  "height":0.4,
  "width": 2.78,
  "horizontal_position":0.31,
  "vertical_position": 1.79,
  "content":[
    {
      "text": "UMK Daerah Tahun 2022",
      "option":{
        "font_size" : 18,
        "bold": True
      }
    }
  ]
}
content3 ={
  "height":0.4,
  "width": 3.75,
  "horizontal_position":0.31,
  "vertical_position": 3.67,
  "content":[
    {
      "text": "Pergeseran UMK Daerah",
      "option":{
        "font_size" : 18,
        "bold": True
      }
    }
  ]
}
chart1 = {
  "height":2.91,
  "width": 10.22,
  "horizontal_position":0.31,
  "vertical_position": 4.08,
  "widget":{
    "type": "chart",
    "value": {
      "chart_type" : "line",
      "chart_data":{
        "catagories_name": ["2020","2021","2022"],
        "series":[
          {
            "series_name":"Kabupaten Lebak",
            "data": [2710654, 2751313, 2773590]
          },
          {
            "series_name" : "Kabupaten Pandeglang",
            "data": [2758909 , 2800292, 2800292]
          },
          {
            "series_name" : "Kabupaten Serang",
            "data": [4152887 , 4251180 , 4125186]
          },
          {
            "series_name" : "Kabupaten Tangerang",
            "data": [4168268 ,4230792 ,4230792]
          },
          {
            "series_name" : "Kota Cilegon",
            "data": [4246081 ,4309772 , 4430254]
          },
          {
            "series_name" : "Kota Serang",
            "data": [3773940 , 3830549 , 3850526]
          },
          {
            "series_name" : "Kota Tanggerang",
            "data": [4119029 , 4262015 , 4285798]
          },
          {
            "series_name" : "Kota Tanggerang Selatan",
            "data": [4168268 , 4230792 , 4280214]
          }
        ]
      },
      "options":{
        "line_type":"LINE_MARKERS",
          "legend":{
              "legend_potision": "right",
              "include_in_layout": False,
              "font_style_setting":{
              "font_size": 14,
              "bold": False
          }
          },
          "value_axis":{
              "minimum_scale": 2500000,
              "maximum_scale": 5000000,
              "format_number": "number",
              "major_gridlines": True
          },
          "category_axis":{
              "major_gridlines": True
          },
          "series":{
              "color": ["5B9BD5","ED7D31","A5A5A5","FFC000","5B9BD5","70AD47","264478","9E480E"] 
          }
    }
  }
}}

# Slide 1
left , top , width , height = Inches(0.3) , Inches(3.62) , Inches(10.22) ,Inches(3.41)
rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
rect.fill.solid()
rect.fill.fore_color.rgb = RGBColor(255, 255, 255) # warna putih
rect.line.color.rgb = RGBColor(255, 255, 255) # warna border sama dengan warna latar

horizontal_position, vertical_position, picture_height , picture_widht ,  = Inches(0.16) , Inches(1.02) , Inches(0.4) ,Inches(4.23)
pic = slide.shapes.add_picture("gambar1.png", horizontal_position, vertical_position,picture_widht,picture_height)

create_content(content=content1)
create_content(content=content2)
create_content(content=content3)
create_chart(widget=chart1)


# Simpan perubahan
prs.save("example_updated.pptx")
os.startfile('example_updated.pptx')
