import collections 
import collections.abc
from pptx import Presentation
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Pt,Inches
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.chart.data import XySeriesData, XyChartData
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_LABEL_POSITION
import os
import json
import traceback
from setting_option import chart_option

def ukuran(data):
        try:
            return round((data / 96),2)

        except Exception:
            traceback.print_exc()

def toIch(cm):
        try:
            return cm * 0.39

        except Exception:
            traceback.print_exc()

def generate_chart(config:Presentation):

    prs = Presentation()
    prs.slide_width = Inches(int(ukuran(config.get('paperSize')[0])))
    prs.slide_height = Inches(int(ukuran(config.get('paperSize')[1])))

    slides = config.get('slides')

    # print(json.dumps(slides, indent=4))

    for slide in slides:
        slide_ppt = prs.slides.add_slide(prs.slide_layouts[6])
        for widget in slide["widgets"]:
            # print(json.dumps(widget, indent=4))
            widget_attr = widget['widget']

            # chart Condition
            if widget_attr['type'] == 'chart':
                match widget_attr['value']['chart_type']:
                     
                    #  Pie Chart Condition
                     case "pie":
                        #   print("ini Pie")
                          chart_atribute = widget_attr['value']["chart_data"]

                          chart_data = ChartData()   
                          chart_data.categories = chart_atribute["catagories_name"]

                          series = chart_atribute["series"]
                        #   print(series)
                          chart_data.add_series(series["series_name"], series["category_data"])
                          
                          x, y, cx, cy = Inches(toIch(widget['left'])), Inches(toIch(widget['top'])), Inches(toIch(widget['width'])), Inches(toIch(widget['height']))
                        #   chart = charts.chart
                          if widget_attr['value'].get("option"):
                            if widget_attr['value']["option"]["exploded"]:
                                charts = slide_ppt.shapes.add_chart(XL_CHART_TYPE.PIE_EXPLODED, x, y, cx, cy, chart_data)
                                chart_option(charts=charts, option=widget_attr['value']["option"], chartCategory=series)
                            else:
                                charts = slide_ppt.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data)
                                chart_option(charts=charts, option=widget_attr['value']["option"])
                          

    prs.save('test2.pptx')
    os.startfile('test2.pptx')
            

if __name__ == "__main__":
    f = open ('config.json', "r")
    data = json.loads(f.read())
    generate_chart(config=data)