import collections 
import collections.abc
from pptx import Presentation
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Pt,Inches
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.chart.data import XySeriesData, XyChartData

def add_chart(slide, chartName:str):
    textbox = slide.shapes.add_textbox(Inches(3),   #left
                                               Inches(1), #top
                                               Inches(5), #width
                                               Inches(2) )#height

    tf = textbox.text_frame
    tf.word_wrap = True
                
    p = tf.paragraphs[0]
    run = p.add_run()
    run.font.size = Pt(21)
                
    run.font.bold = True

    align = PP_ALIGN.CENTER
                
    p.alignment = align

    run.text = chartName

    if chartName == "COLUMN_CLUSTERED":
        chart_data = ChartData()   
        chart_data.categories = ['Twitter', 'Youtube', 'Instagram', 'Facebook', 'Tiktok']
        chart_data.add_series('25 Jan 00:00', (451, 54, 32, 18, 3))
        chart_data.add_series('25 Jan 01:00', (354, 28, 58, 13, 6))
        chart_data.add_series('25 Jan 02:00', (354, 28, 58, 13, 4))

        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(3)

        chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart

    elif chartName == "COLUMN_STACKED":
        chart_data = ChartData()   
        chart_data.categories = ['Twitter', 'Youtube', 'Instagram', 'Facebook', 'Tiktok']
        chart_data.add_series('25 Jan 00:00', (451, 54, 32, 18, 3))
        chart_data.add_series('25 Jan 01:00', (354, 28, 58, 13, 6))
        chart_data.add_series('25 Jan 02:00', (354, 28, 58, 13, 4))

        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(3)

        chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, x, y, cx, cy, chart_data).chart
    
    elif chartName == "COLUMN_STACKED_100":
        chart_data = ChartData()   
        chart_data.categories = ['Twitter', 'Youtube', 'Instagram', 'Facebook', 'Tiktok']
        chart_data.add_series('25 Jan 00:00', (451, 54, 32, 18, 3))
        chart_data.add_series('25 Jan 01:00', (354, 28, 58, 13, 6))
        chart_data.add_series('25 Jan 02:00', (354, 28, 58, 13, 4))

        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(3)

        chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED_100, x, y, cx, cy, chart_data).chart

    elif chartName == "LINE":
        chart_data = ChartData()   
        chart_data.categories = ['Twitter', 'Youtube', 'Instagram', 'Facebook', 'Tiktok']
        chart_data.add_series('25 Jan 00:00', (451, 54, 32, 18, 3))
        chart_data.add_series('25 Jan 01:00', (354, 28, 58, 13, 6))
        chart_data.add_series('25 Jan 02:00', (354, 28, 58, 13, 4))

        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(3)

        chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data).chart

    elif chartName == "LINE_STACKED":
        chart_data = ChartData()   
        chart_data.categories = ['Twitter', 'Youtube', 'Instagram', 'Facebook', 'Tiktok']
        chart_data.add_series('25 Jan 00:00', (451, 54, 32, 18, 3))
        chart_data.add_series('25 Jan 01:00', (354, 28, 58, 13, 6))
        chart_data.add_series('25 Jan 02:00', (354, 28, 58, 299, 50))

        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(3)

        chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_STACKED, x, y, cx, cy, chart_data).chart

    elif chartName == "LINE_MARKERS_STACKED":
        chart_data = ChartData()   
        chart_data.categories = ['Twitter', 'Youtube', 'Instagram', 'Facebook', 'Tiktok']
        chart_data.add_series('25 Jan 00:00', (451, 54, 32, 18, 3))
        chart_data.add_series('25 Jan 01:00', (354, 28, 58, 13, 6))
        chart_data.add_series('25 Jan 02:00', (354, 28, 58, 13, 4))

        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(3)

        chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS_STACKED, x, y, cx, cy, chart_data).chart

    elif chartName == "LINE_MARKERS_STACKED_100":
        chart_data = ChartData()   
        chart_data.categories = ['Twitter', 'Youtube', 'Instagram', 'Facebook', 'Tiktok']
        chart_data.add_series('25 Jan 00:00', (451, 54, 32, 18, 3))
        chart_data.add_series('25 Jan 01:00', (354, 28, 58, 13, 6))
        chart_data.add_series('25 Jan 02:00', (354, 28, 58, 13, 4))

        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(3)

        chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS_STACKED_100, x, y, cx, cy, chart_data).chart

    elif chartName == "LINE_STACKED_100":
        chart_data = ChartData()   
        chart_data.categories = ['Twitter', 'Youtube', 'Instagram', 'Facebook', 'Tiktok']
        chart_data.add_series('25 Jan 00:00', (451, 54, 32, 18, 3))
        chart_data.add_series('25 Jan 01:00', (354, 28, 58, 13, 6))
        chart_data.add_series('25 Jan 02:00', (354, 28, 58, 13, 4))

        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(3)

        chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_STACKED_100, x, y, cx, cy, chart_data).chart

    elif chartName == "LINE_MARKERS":
        chart_data = ChartData()   
        chart_data.categories = ['Twitter', 'Youtube', 'Instagram', 'Facebook', 'Tiktok']
        chart_data.add_series('25 Jan 00:00', (451, 54, 32, 18, 3))
        chart_data.add_series('25 Jan 01:00', (354, 28, 58, 13, 6))
        chart_data.add_series('25 Jan 02:00', (354, 28, 58, 13, 4))

        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(3)

        chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, x, y, cx, cy, chart_data).chart

    elif chartName == "PIE":
        chart_data = ChartData()   
        chart_data.categories = ['Twitter', 'Youtube', 'Instagram', 'Facebook', 'Tiktok']
        chart_data.add_series('25 Jan 00:00', (451, 54, 32, 18, 3))

        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(3)

        chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data).chart
        chart.has_legend = True

    elif chartName == "PIE_EXPLODED":
        chart_data = ChartData()   
        chart_data.categories = ['Twitter', 'Youtube', 'Instagram', 'Facebook', 'Tiktok']
        chart_data.add_series('25 Jan 00:00', (451, 54, 32, 18, 3))

        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(3)

        chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE_EXPLODED, x, y, cx, cy, chart_data).chart
        chart.has_legend = True

    elif chartName == "DOUGHNUT":
        chart_data = ChartData()   
        chart_data.categories = ['Twitter', 'Youtube', 'Instagram', 'Facebook', 'Tiktok']
        chart_data.add_series('25 Jan 00:00', (451, 54, 32, 18, 3))

        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(3)

        chart = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, chart_data).chart
        chart.has_legend = True

    elif chartName == "DOUGHNUT_EXPLODED":
        chart_data = ChartData()   
        chart_data.categories = ['Twitter', 'Youtube', 'Instagram', 'Facebook', 'Tiktok']
        chart_data.add_series('25 Jan 00:00', (451, 54, 32, 18, 3))

        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(3)

        chart = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT_EXPLODED, x, y, cx, cy, chart_data).chart
        chart.has_legend = True

    elif chartName == "BAR_CLUSTERED":
        chart_data = ChartData()   
        chart_data.categories = ['Twitter', 'Youtube', 'Instagram', 'Facebook', 'Tiktok']
        chart_data.add_series('25 Jan 00:00', (451, 54, 32, 18, 3))
        chart_data.add_series('25 Jan 00:00', (451, 54, 32, 18, 3))
        chart_data.add_series('25 Jan 00:00', (451, 54, 32, 18, 3))

        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(3)

        chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data).chart
        # chart.has_legend = True
    
    elif chartName == "BAR_STACKED":
        chart_data = ChartData()   
        chart_data.categories = ['Twitter', 'Youtube', 'Instagram', 'Facebook', 'Tiktok']
        chart_data.add_series('25 Jan 00:00', (451, 54, 32, 18, 3))
        chart_data.add_series('25 Jan 00:00', (451, 54, 32, 18, 3))
        chart_data.add_series('25 Jan 00:00', (451, 54, 32, 18, 3))

        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(3)

        chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED, x, y, cx, cy, chart_data).chart
        # chart.has_legend = True
    
    elif chartName == "BAR_STACKED_100":
        chart_data = ChartData()   
        chart_data.categories = ['Twitter', 'Youtube', 'Instagram', 'Facebook', 'Tiktok']
        chart_data.add_series('25 Jan 00:00', (451, 54, 32, 18, 3))
        chart_data.add_series('25 Jan 00:00', (451, 54, 32, 18, 3))
        chart_data.add_series('25 Jan 00:00', (451, 54, 32, 18, 3))

        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(3)

        chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED_100, x, y, cx, cy, chart_data).chart
        # chart.has_legend = True
    
    elif chartName == "AREA":
        chart_data = ChartData()   
        chart_data.categories = ['Twitter', 'Youtube', 'Instagram', 'Facebook', 'Tiktok']
        chart_data.add_series('25 Jan 00:00', (451, 54, 32, 18, 3))
        chart_data.add_series('25 Jan 00:00', (122, 31, 123, 3, 3))
        chart_data.add_series('25 Jan 00:00', (344, 4, 45, 1, 3))

        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(3)

        chart = slide.shapes.add_chart(XL_CHART_TYPE.AREA, x, y, cx, cy, chart_data).chart
        # chart.has_legend = True
    
    elif chartName == "AREA_STACKED":
        chart_data = ChartData()   
        chart_data.categories = ['Twitter', 'Youtube', 'Instagram', 'Facebook', 'Tiktok']
        chart_data.add_series('25 Jan 00:00', (451, 54, 32, 18, 3))
        chart_data.add_series('25 Jan 00:00', (122, 31, 123, 3, 3))
        chart_data.add_series('25 Jan 00:00', (344, 4, 45, 1, 3))

        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(3)

        chart = slide.shapes.add_chart(XL_CHART_TYPE.AREA_STACKED, x, y, cx, cy, chart_data).chart
        # chart.has_legend = True


    elif chartName == "AREA_STACKED_100":
        chart_data = ChartData()   
        chart_data.categories = ['Twitter', 'Youtube', 'Instagram', 'Facebook', 'Tiktok']
        chart_data.add_series('25 Jan 00:00', (451, 54, 32, 18, 3))
        chart_data.add_series('25 Jan 00:00', (122, 31, 123, 3, 3))
        chart_data.add_series('25 Jan 00:00', (344, 4, 45, 1, 3))

        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(3)

        chart = slide.shapes.add_chart(XL_CHART_TYPE.AREA_STACKED_100, x, y, cx, cy, chart_data).chart
        # chart.has_legend = True
