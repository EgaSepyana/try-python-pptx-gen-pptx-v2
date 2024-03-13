import collections 
import collections.abc
from pptx import Presentation
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Pt,Inches
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.chart.data import XySeriesData, XyChartData
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.chart import XL_TICK_MARK
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_MARKER_STYLE
import os
from pptx.enum.chart import XL_LABEL_POSITION

def pieChart():
    prs = Presentation()

    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    chart_data = ChartData()
    categories = ["2020", "2021", "2022"] 
    chart_data.categories = categories
    chart_data.add_series('Kabupaten Lebak', (2710654, 2751313, 2773590 ))
    chart_data.add_series('Kabupaten Pandeglang', (2758909 , 2800292, 2800292))
    chart_data.add_series('Kabupaten Serang', ( 4152887 , 4251180 , 4125186))
    chart_data.add_series('Kabupaten Tangerang', ( 4168268 ,4230792 ,4230792))
    chart_data.add_series('Kota Cilegon', ( 4246081 ,4309772 , 4430254 ))
    chart_data.add_series('Kota Serang', ( 3773940 , 3830549 , 3850526  ))
    chart_data.add_series('Kota Tanggerang', ( 4119029 , 4262015 , 4285798))
    chart_data.add_series('Kota Tanggerang Selatan', ( 4168268 , 4230792 , 4280214 ))
    # chart_data.categories = categories
    # series1 = chart_data.add_series("Kabupaten Lebak")
    # series1_data = [ 2710654 , 2751313 ,2773590 ]
    # series2 = chart_data.add_series("Kabupaten Pandeglang")
    # series2_data = [ 2758909 ,2800292 , 2800292 ]
    # series3 = chart_data.add_series("Kabupaten Serang")
    # series3_data = [ 4152887 , 4251180 , 4125186]
    # series4 = chart_data.add_series("Kabupaten Tangerang")
    # series4_data = [ 4168268 ,4230792 ,4230792]
    # series5 = chart_data.add_series("Kota Cilegon")
    # series5_data = [  4246081 ,4309772 , 4430254 ]
    # series6 = chart_data.add_series("Kota Serang")
    # series6_data = [  3773940 , 3830549 , 3850526  ]

#     Time	Kabupaten Lebak	 Kabupaten Pandeglang	Kabupaten Serang	Kabupaten Tangerang	Kota Cilegon	Kota Serang	Kota Tangerang	Kota Tangerang Selatan
    # 2019								
    # 2020	 2.710.654 	 2.758.909 	 4.152.887 	 4.168.268 	 4.246.081 	 3.773.940 	 4.119.029 	 4.168.268 
    # 2021	 2.751.313 	 2.800.292 	 4.251.180 	 4.230.792 	 4.309.772 	 3.830.549 	 4.262.015 	 4.230.792 
    # 2022	 2.773.590 	 2.800.292 	 4.125.186 	 4.230.792 	 4.430.254 	 3.850.526 	 4.285.798 	 4.280.214 


    # for i in range(len(categories)):
    #     series1.add_data_point(x=categories[i], y=series1_data[i])
    
    # for i in range(len(categories)):
    #     series2.add_data_point(x=categories[i], y=series2_data[i])
    # for i in range(len(categories)):
    #     series3.add_data_point(x=categories[i], y=series3_data[i])
    # for i in range(len(categories)):
    #     series4.add_data_point(x=categories[i], y=series4_data[i])
    
    # for i in range(len(categories)):
    #     series5.add_data_point(x=categories[i], y=series5_data[i])

    # for i in range(len(categories)):
    #     series6.add_data_point(x=categories[i], y=series6_data[i])

    x, y, cx, cy = Inches(0.44), Inches(2.08), Inches(9.49), Inches(2.91)

    charts = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, x, y, cx, cy, chart_data)
    chart = charts.chart

    chart.has_title = False
    
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(12)
    
    # Value Axis
    axis = chart.value_axis
    axis.minimum_scale = 2500000
    axis.maximum_scale = 5000000
    axis.tick_labels.number_format = '#,###'
    # axis.has_minor_gridlines = True
    axis.has_major_gridlines = True
    # chart.value_axis.minor_tick_mark = XL_TICK_MARK.INSIDE

    # Category Axis
    category_axis = chart.category_axis
    category_axis.has_minor_gridlines = True



    # Chart 2 

    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    chart_data = ChartData()
    categories = ["Negri", "Swasta"]
    chart_data.categories = categories
    chart_data.add_series('Sekolah', (245, 52 ))

    x, y, cx, cy = Inches(-0.09), Inches(1.93), Inches(3.42), Inches(3.42)

    charts = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, chart_data)
    chart = charts.chart

    chart.has_title = True
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM

    data_label = chart.plots[0]
    data_label.has_data_labels = True
    data_labels = data_label.data_labels
    data_labels.show_percentage = True
    data_labels.font.size = Pt(14)
    # data_labels.font.bold = True

    plotstyle = {
        "color" : ["F8AB13", "1DB7D9"]
    }
    for idx, point in enumerate(chart.series[0].points):
                    col_idx = idx % len(plotstyle["color"])           
                    point.format.fill.solid()
                    point.format.fill.fore_color.rgb = RGBColor.from_string(plotstyle["color"][col_idx])
    

    chart_data = ChartData()
    categories = ["Negri", "Swasta"]
    chart_data.categories = categories
    chart_data.add_series('Guru', (3221, 505 ))

    x, y, cx, cy = Inches(3.27), Inches(1.93), Inches(3.42), Inches(3.42)

    charts = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, chart_data)
    chart = charts.chart

    chart.has_title = True
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM

    data_label = chart.plots[0]
    data_label.has_data_labels = True
    data_labels = data_label.data_labels
    data_labels.show_percentage = True
    data_labels.font.size = Pt(14)
    # data_labels.font.bold = True

    plotstyle = {
        "color" : ["F8AB13", "1DB7D9"]
    }
    for idx, point in enumerate(chart.series[0].points):
                    col_idx = idx % len(plotstyle["color"])           
                    point.format.fill.solid()
                    point.format.fill.fore_color.rgb = RGBColor.from_string(plotstyle["color"][col_idx])


    chart_data = ChartData()
    categories = ["Negri", "Swasta"]
    chart_data.categories = categories
    chart_data.add_series('Murid', (21034, 2451 ))

    x, y, cx, cy = Inches(6.59), Inches(1.93), Inches(3.42), Inches(3.42)

    charts = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT_EXPLODED, x, y, cx, cy, chart_data)
    chart = charts.chart

    chart.has_title = True
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM

    data_label = chart.plots[0]
    data_label.has_data_labels = True
    data_labels = data_label.data_labels
    data_labels.show_percentage = True
    data_labels.font.size = Pt(14)
    # data_labels.font.bold = True

    plotstyle = {
        "color" : ["F8AB13", "1DB7D9"]
    }
    for idx, point in enumerate(chart.series[0].points):
                    col_idx = idx % len(plotstyle["color"])           
                    point.format.fill.solid()
                    point.format.fill.fore_color.rgb = RGBColor.from_string(plotstyle["color"][col_idx])




    # chart 3

    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    chart_data = ChartData()
    categories = ["Kabupaten Badung", "Kabupaten Bangli","Kabupaten Buleleng","Kabupaten Gianyar","Kabupaten Jembrana","Kabupaten Karangasem","Kabupaten Klungkung","Kabupaten Tabanan","Kota Denpasar"] 
    chart_data.categories = categories
    chart_data.add_series('Bekerja', (72.69,82.2,75.07,71.27,75.92,80.75,75.27,75.47,70.91))
    chart_data.add_series('Kabupaten Pandeglang', (6.92,1.86,5.19,7.53,4.52,2.42,5.42,4.21,7.62))

    x, y, cx, cy = Inches(0), Inches(2.08), Inches(10.00), Inches(4.89)

    charts = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED, x, y, cx, cy, chart_data)
    chart = charts.chart

    chart.has_title = False
    
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(12)

    plotstyle = {
        "color" : ["4472C4","ED7D31"]
    }

    for idx, point in enumerate(chart.series):
                    col_idx = idx % len(plotstyle["color"])           
                    point.format.fill.solid()
                    point.format.fill.fore_color.rgb = RGBColor.from_string(plotstyle["color"][col_idx])
    

    # chart 4

    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    chart_data = ChartData()
    categories = ["2018", "2019", "2020"] 
    chart_data.categories = categories
    chart_data.add_series('Kabupaten Tanggerang', ( 850184 , 917146 , 969308 ))
    chart_data.add_series('Kota Tanggerang', (  342045 , 358534 , 371871 ))
    chart_data.add_series('Kota Tanggerang Selatan', (  494850 , 521424 , 544979 ))

    x, y, cx, cy = Inches(0), Inches(2.64), Inches(10), Inches(4.71)

    charts = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, x, y, cx, cy, chart_data)
    chart = charts.chart

    chart.has_title = False
    
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(12)
    
    # Value Axis
    axis = chart.value_axis
    # axis.minimum_scale = 2500000
    # axis.maximum_scale = 5000000
    axis.tick_labels.number_format = '#,###'
    # axis.has_minor_gridlines = True
    axis.has_major_gridlines = True
    # chart.value_axis.minor_tick_mark = XL_TICK_MARK.INSIDE

    # Category Axis
    category_axis = chart.category_axis
    category_axis.has_minor_gridlines = True

    plotstyle = {
        "color" : ["4472C4", "A5A5A5","ED7D31"]
    }

    plot = chart.plots[0]
    series = plot.series

    for i in range(len(series)):
        line = series[i].format.line
        series[i].marker.format.fill.solid()
        series[i].marker.format.fill.fore_color.rgb = RGBColor.from_string(plotstyle["color"][i])
        series[i].marker.style = XL_MARKER_STYLE.CIRCLE
        line.color.rgb = RGBColor.from_string(plotstyle["color"][i])

    data_label = chart.plots[0]
    data_label.has_data_labels = True
    data_labels = data_label.data_labels
    data_labels.show_value = True
    data_labels.number_format = '#,###'
    data_labels.position = XL_LABEL_POSITION.ABOVE
    data_labels.font.size = Pt(16)


    # chart 4

    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    chart_data = ChartData()
    categories = ["Berau", "Kutai Bara","Kutai Kartanegara","Kutai Timur","Paser","Penajam Paser Utara","Kota Balikpapan","Kota Bontang","Kota Samarinda"] 
    chart_data.categories = categories
    chart_data.add_series('2018', ( 149,94,525,181,248,129,518,146,579 ))
    chart_data.add_series('2019', (  181,96,598,209,267,150,564,156,625))

    x, y, cx, cy = Inches(0), Inches(2.64), Inches(10), Inches(4.71)

    charts = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)
    chart = charts.chart

    chart.has_title = False
    
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.TOP
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(12)

    axis = chart.value_axis
    axis.has_major_gridlines = False
    axis.visible = False

    category_axis = chart.category_axis
    category_axis.has_major_gridlines = True
    category_axis.tick_labels.font.size = Pt(10.6)
    category_axis.tick_labels.font.name = "Comic Sans MS"

    data_label = chart.plots[0]
    data_label.has_data_labels = True
    data_labels = data_label.data_labels
    data_labels.show_percentage = True
    data_labels.font.size = Pt(10.6)

    data_label.gap_width = 274
    data_label.overlap = -81
    
    
    plotstyle = {
        "color" : ["4472C4","ED7D31"]
    }

    for idx, point in enumerate(chart.series):
                    col_idx = idx % len(plotstyle["color"])           
                    point.format.fill.solid()
                    point.format.fill.fore_color.rgb = RGBColor.from_string(plotstyle["color"][col_idx])
    # for idx, series in enumerate(chart.series):
    #                 col_idx = idx % len(plotstyle["color"])           
    #                 series.format.fill.solid()
    #                 series.format.fill.fore_color.rgb = RGBColor.from_string(plotstyle["color"][col_idx])

    # Chart 5

    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    chart_data = ChartData()
    categories = ["MAJENE", "MAMASA", "MAMUJU","MAMUJU TENGAH","MAMUJU UTARA","POLEWALI MANDAR"] 
    chart_data.categories = categories
    chart_data.add_series('BERKEMBANG', ( 40,18,48,32,32,94 ))
    chart_data.add_series('MAJU', ( 2,1,7,3,10,12 ))
    chart_data.add_series('MANDIRI', ( None,None,None,1,None,None ))
    chart_data.add_series('SANGAT TERTINGGAL', ( 6,30,None,None,1,None ))
    chart_data.add_series('TERTINGGAL', ( 14,119,33,18,16,38 ))

    x, y, cx, cy = Inches(0), Inches(2.64), Inches(10), Inches(4.71)

    charts = slide.shapes.add_chart(XL_CHART_TYPE.RADAR, x, y, cx, cy, chart_data)
    chart = charts.chart

    chart.has_title = False

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.TOP
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(12)

    axis = chart.value_axis
    axis.minimum_scale = 0
    axis.maximum_scale = 150
    # axis.tick_labels.number_format = '#,###'
    axis.visible = False
    # axis.has_minor_gridlines = True
    # axis.has_major_gridlines = True
    # chart.value_axis.minor_tick_mark = XL_TICK_MARK.INSIDE

    # Category Axis
    # category_axis = chart.category_axis
    # category_axis.has_minor_gridlines = True

    plotstyle = {
        "color" : ["4472C4","ED7D31","A5A5A5","FFC000","5B9BD5"]
    }

    for i in range(len(chart.plots[0].series)):
        line = chart.plots[0].series[i].format.line
        line.color.rgb = RGBColor.from_string(plotstyle["color"][i])


    # chart 6
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # menentukan posisi dan ukuran grafik
    x, y, cx, cy = Inches(0), Inches(2.58), Inches(4.82), Inches(4.94)

    chart_data = ChartData()
    categories = ["2018", "2019", "2020","2021"] 
    chart_data.categories = categories
    chart_data.add_series('Kota Surabaya', ( 0.358,	0.406,	0.344,	0.351 ))
    chart_data.add_series('Siduarjo', ( 0.346,	0.311	,0.34,	0.347 ))

    charts = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)
    chart = charts.chart

    data_label = chart.plots[0]
    data_label.has_data_labels = True
    data_labels = data_label.data_labels
    data_labels.show_percentage = True
    data_labels.font.size = Pt(12)

    data_label.gap_width = 100
    data_label.overlap = -24

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(12)

    axis = chart.value_axis
    axis.has_major_gridlines = False
    axis.visible = False

    category_axis = chart.category_axis
    category_axis.has_major_gridlines = False
    category_axis.tick_labels.font.size = Pt(12)

    plotstyle = {
        "color" : ["B4D92A","1DB7D9"]
    }

    for idx, point in enumerate(chart.series):
                    col_idx = idx % len(plotstyle["color"])           
                    point.format.fill.solid()
                    point.format.fill.fore_color.rgb = RGBColor.from_string(plotstyle["color"][col_idx])

    chart_data = ChartData()
    categories = ["2014", "2015", "2016","2017","2018","2019","2020","2021"] 
    chart_data.categories = categories
    chart_data.add_series('Kota Surabaya', ( 73.85,73.85,73.87,73.88,73.98,74.13,74.18,74.18 ))
    chart_data.add_series('Siduarjo', (  73.43,73.63,73.67,73.71,73.82,73.98,74.04,74.06 ))

    x, y, cx, cy = Inches(4.71), Inches(2.7), Inches(5.29), Inches(4.82)

    charts = slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data)
    chart = charts.chart

    axis = chart.value_axis
    axis.has_major_gridlines = False
    axis.visible = False

    category_axis = chart.category_axis
    category_axis.has_major_gridlines = False
    category_axis.tick_labels.font.size = Pt(12)

    data_label = chart.plots[0]
    data_label.has_data_labels = True
    data_labels = data_label.data_labels
    data_labels.show_percentage = True
    data_labels.font.size = Pt(12)

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(12)

    plot = chart.plots[0]

    for i in plot.series:
        i.smooth = True

    plotstyle = {
        "color" : ["B4D92A","1DB7D9"]
    }

    for i in range(len(chart.plots[0].series)):
        line = chart.plots[0].series[i].format.line
        line.color.rgb = RGBColor.from_string(plotstyle["color"][i])


    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # menentukan posisi dan ukuran grafik
    x, y, cx, cy = Inches(0), Inches(3.23), Inches(4.51), Inches(4.27)

    chart_data = ChartData()
    categories = ["Kota Blitar", "Kota Pasuruan","Kota Probolinggo","Kota Malang","Sidoarjo","Kota Kediri","Kota Surabaya","Nganjuk"]
    chart_data.categories = categories
    chart_data.add_series('2021', (97.36, 96.27,96.01,94.42,94.26,94.21,93.9,93.76) )

    charts = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)
    chart = charts.chart

    plotstyle = {
        "color" : ["4472C4"]
    }


    for idx, point in enumerate(chart.series):
                    col_idx = idx % len(plotstyle["color"])           
                    point.format.fill.solid()
                    point.format.fill.fore_color.rgb = RGBColor.from_string(plotstyle["color"][col_idx])

    x, y, cx, cy = Inches(4.35), Inches(2.32), Inches(5.65), Inches(5.09)

    chart_data = ChartData()
    categories = ["Magetan","Nganjuk","Banyuwangi","Kota Probolinggo","Kota Kediri","Mojokerto","Kota Batu","Gresik","Kota Pasuruan","Kota Mojokerto","Kota Blitar","Sidoarjo","Kota Madiun","Kota Malang","Kota Surabaya"]
    chart_data.categories = categories
    chart_data.add_series('2021', (11833,12172,12217,12245,12359,12844,12887,13280,13354,13610,13816,14578,16095,16663,17862) )

    charts = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data)
    chart = charts.chart

    plotstyle = {
        "color" : ["4472C4"]
    }

    for idx, point in enumerate(chart.series):
                    col_idx = idx % len(plotstyle["color"])           
                    point.format.fill.solid()
                    point.format.fill.fore_color.rgb = RGBColor.from_string(plotstyle["color"][col_idx])

    axis = chart.value_axis
    axis.has_major_gridlines = False
    axis.minimum_scale = 0
    axis.maximum_scale = 25000
    axis.visible = False

    category_axis = chart.category_axis
    category_axis.has_major_gridlines = False
    category_axis.tick_labels.font.size = Pt(16)

    

    data_label = chart.plots[0]
    data_label.has_data_labels = True
    data_labels = data_label.data_labels
    data_labels.show_percentage = True
    data_labels.font.size = Pt(16)
    
    # data_label.gap_width = 100
    

    prs.save('pieChartEx.pptx')
    os.startfile('pieChartEx.pptx')

if __name__ == "__main__":
    pieChart()