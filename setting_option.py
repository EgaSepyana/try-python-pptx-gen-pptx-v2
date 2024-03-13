import collections 
import collections.abc
from pptx import Presentation
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Pt,Inches
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.chart.data import XySeriesData, XyChartData
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_LABEL_POSITION
import os
import json
import traceback


def chart_option(charts, option:dict,chartCategory:dict):
    print(json.dumps(option, indent=4))
    chart = charts.chart
    chart.has_title = False
    if option.get("title"):
        chart.has_title = True
        if option["title"].get("font_style_setting"):
            font_style_setting = option["title"]["font_style_setting"]
            chart.chart_title.has_text_frame = True
            title_text_frame = chart.chart_title.text_frame
            title = title_text_frame.paragraphs[0].add_run()
            edit_text(title, font_style_setting)

    if option.get("plots"):
        plots = option["plots"]

        if plots.get('style'):
            plotstyle = plots["style"]
            if plotstyle["color"] and (len(plotstyle["color"]) == len(chartCategory["category_data"])):
                for idx, point in enumerate(chart.series[0].points):
                    col_idx = idx % len(plotstyle["color"])           
                    point.format.fill.solid()
                    point.format.fill.fore_color.rgb = RGBColor.from_string(plotstyle["color"][col_idx])

        title.text = "Judul"

def edit_text(textframe, text_config):

    if text_config.get("size_pt"):
        textframe.font.size = Pt(int(text_config["size_pt"]))
    else:
        textframe.font.size = Pt(30)

    if text_config["bold"] == True:
        textframe.font.bold = True

    if text_config["underline"] == True:
        textframe.font.underline = True

    if text_config["italic"] == True:
        textframe.font.italic = True

    if text_config.get("font_color_hex"):
        textframe.font.color.rgb = RGBColor.from_string(text_config["font_color_hex"])
    