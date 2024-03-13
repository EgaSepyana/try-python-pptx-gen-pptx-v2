import collections 
import collections.abc
from pptx import Presentation
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Pt,Inches
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.chart.data import XySeriesData, XyChartData
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.chart import XL_TICK_MARK
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_MARKER_STYLE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import os
from pptx.enum.chart import XL_LABEL_POSITION

prs = Presentation("Contoh.pptx")

blank_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(blank_slide_layout)


left , top , width , height = Inches(9.43) , Inches(0.21) , Inches(3.63) ,Inches(7.28)
rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
rect.fill.solid()
rect.fill.fore_color.rgb = RGBColor.from_string("EBB55A") # warna putih
rect.fill.fore_color.brightness = 0.71
rect.line.color.rgb = RGBColor(255, 255, 255) # warna border sama dengan warna latar


left , top , width , height = Inches(8.19) , Inches(0.28) , Inches(4.88) ,Inches(0.36)
rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
rect.fill.solid()
rect.fill.fore_color.rgb = RGBColor.from_string("EBB55A") # warna putih
rect.line.color.rgb = RGBColor.from_string("EBB55A") # warna border sama dengan warna latar

horizontal_position , vertical_position , width , height = Inches(1.35),Inches(0.24),Inches(7.82),Inches(0.78)

txBox = slide.shapes.add_textbox(horizontal_position, vertical_position, width, height)

text_frame = txBox.text_frame
text_frame.word_wrap = True

p = text_frame.paragraphs[0]

run1 = p.add_run()
run1.font.size = Pt(36)
p.alignment = PP_ALIGN.LEFT
run1.text = "Exposure"
run2 = p.add_run()
run2.font.size = Pt(36)
run2.font.bold = True
run2.text = " \"Penolakan Vaksin\""
run1.font.color.rgb = RGBColor.from_string("7F7F7F")
run2.font.color.rgb = RGBColor.from_string("7F7F7F")

horizontal_position , vertical_position , width , height = Inches(4.17),Inches(0.91),Inches(3.84),Inches(0.4)
txBox = slide.shapes.add_textbox(horizontal_position, vertical_position, width, height)
text_frame = txBox.text_frame
text_frame.word_wrap = True

p2 = text_frame.paragraphs[0]
run2 = p2.add_run()
run2.font.size = Pt(18)
run2.text = " 1 Januari 2021 – 31 Maret 2021"
run2.font.color.rgb = RGBColor.from_string("D84E2E")

horizontal_position , vertical_position , width , height = Inches(0.38),Inches(5.14),Inches(8.71),Inches(2.06)
txBox = slide.shapes.add_textbox(horizontal_position, vertical_position, width, height)
text_frame = txBox.text_frame
text_frame.word_wrap = True

p2 = text_frame.paragraphs[0]
run2 = p2.add_run()
run2.font.size = Pt(14)
run2.text = "Pembahasan Vaksin mulai ramai dibicarakan netizen pada bulan Januari (sejak wacana vaksin dimulai),"
run2.font.bold = True

run2 = p2.add_run()
run2.font.size = Pt(14)
run2.text = " di mana pemerintah mengharuskan seluruh masyarakat untuk mengikuti vaksinasi, namun diikuti dengan pernyataan dari"

run2 = p2.add_run()
run2.font.size = Pt(14)
run2.text = " Anggota DPR Ribka Tjiptaning dalam forum resmi legislatif yang menyatakan menolak menerima vaksin corona buatan perusahaan farmasi asal China, Sinovac."
run2.font.bold = True


run2 = p2.add_run()
run2.font.size = Pt(14)
run2.text = "Vaksinasi Covid-19 tahap pertama mulai bergulir di berbagai daerah walau sejumlah kalangan masih enggan mengikuti salah satu upaya mengatasi pandemi covid-19."

run2 = p2.add_run()
run2.font.size = Pt(14)
run2.text = "Trend pembicaraan penolakan vaksin terus bergerak menurun. Sementara tagar #TolakDivaksinSinovac sempat mencuat di Twitter karena dicuitkan belasan ribu kali."
run2.font.bold = True

chart_data = ChartData()   
chart_data.categories = ['Jumlah reaksi pembicaraan vaksin secara keseluruhan', 'jumlah reaksi penolakan vaksin']
chart_data.add_series('25 Jan 00:00', [83.797,34.854])

x, y, cx, cy = Inches(9.41), Inches(0.91), Inches(3.66), Inches(2.97)

charts = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data)
chart = charts.chart
chart.has_title = False
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False
chart.legend.font.size = Pt(12)
plotstyle = {
    "color" : ["D84E2E", "7E979B"]
}
for idx, point in enumerate(chart.series[0].points):
                col_idx = idx % len(plotstyle["color"])           
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = RGBColor.from_string(plotstyle["color"][col_idx])
data_label = chart.plots[0]
data_label.has_data_labels = True
data_labels = data_label.data_labels
data_labels.position = XL_LABEL_POSITION.INSIDE_END
data_labels.show_percentage = True
data_labels.show_value = False
data_labels.show_category_name = False
data_labels.font.size = Pt(16)
data_labels.font.bold = True
data_labels.font.color.rgb = RGBColor.from_string("FFFFFF")

a = """
    x: 9.31,
      y: 4.63, 
      w: 3.66, 
      h: 2.5,
     """

chart_data = ChartData()   
chart_data.categories = ['Jumlah reaksi pembicaraan vaksin secara keseluruhan', 'jumlah reaksi penolakan vaksin']
chart_data.add_series('25 Jan 00:00', [83.797,34.854])

x, y, cx, cy = Inches(9.31), Inches(4.63), Inches(3.66), Inches(2.5)

charts = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data)
chart = charts.chart
chart.has_title = False
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False
chart.legend.font.size = Pt(9)
plotstyle = {
    "color" : ["D84E2E", "7E979B"]
}
for idx, point in enumerate(chart.series[0].points):
                col_idx = idx % len(plotstyle["color"])           
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = RGBColor.from_string(plotstyle["color"][col_idx])
data_label = chart.plots[0]
data_label.has_data_labels = True
data_labels = data_label.data_labels
data_labels.position = XL_LABEL_POSITION.INSIDE_END
data_labels.show_percentage = True
data_labels.show_value = False
data_labels.show_category_name = False
data_labels.font.size = Pt(16)
data_labels.font.bold = True
data_labels.font.color.rgb = RGBColor.from_string("FFFFFF")


"""{
        type: pptx.charts.BAR,
        data: [{
          name: "Post",
          labels: ["Jan-21","Feb-21","Mar-21"],
          values: [29173,3534,3314],
        }
      ],
        options:{
          showLabel: true,
          chartColors: ["D84E2E"],
          dataLabelColor: "000000",
          dataLabelFontFace: "Arial",
          dataLabelFontSize: 9,
          dataLabelPosition: "outEnd",
          showValue: true,
        }
      },"""
chart_data = ChartData()

categories = ["Jan-21","Feb-21","Mar-21"]

chart_data.categories = categories

chart_data.add_series('Post', [29173,3534,3314])

x, y, cx, cy = Inches(0.32), Inches(1.3), Inches(8.51), Inches(3.61)

charts = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)
chart = charts.chart

data_label = chart.plots[0]
data_label.has_data_labels = True
data_labels = data_label.data_labels
data_labels.show_percentage = True
data_labels.font.size = Pt(12)
data_label.gap_width = 100
data_label.overlap = -24

chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.LEFT
chart.legend.include_in_layout = False
chart.legend.font.size = Pt(12)

axis = chart.value_axis
axis.has_major_gridlines = True
axis.visible = True

category_axis = chart.category_axis
category_axis.has_major_gridlines = False
category_axis.tick_labels.font.size = Pt(12)

plotstyle = {
        "color" : ["D84E2E"]
    }

for idx, point in enumerate(chart.series):
             col_idx = idx % len(plotstyle["color"])           
             point.format.fill.solid()
             point.format.fill.fore_color.rgb = RGBColor.from_string(plotstyle["color"][col_idx])

prs.save('Contoh2.pptx')
os.startfile('Contoh2.pptx')