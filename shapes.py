import collections 
import collections.abc
from pptx import Presentation
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Pt,Inches
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.chart.data import XySeriesData, XyChartData
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.chart import XL_TICK_MARK
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_MARKER_STYLE
from pptx.enum.text import PP_ALIGN
import os
from pptx.enum.chart import XL_LABEL_POSITION

def add_shapes():
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.dml.color import RGBColor
        from pptx.enum.dml import MSO_SHADOW_STYLE
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        left = top = width = height = Inches(1)
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(255, 255, 255) # warna putih
        rect.line.color.rgb = RGBColor(255, 255, 255) # warna border sama dengan warna latar

        prs.save('rectangle_shape_with_shadow.pptx')

if __name__ == "__main__":
    add_shapes()
