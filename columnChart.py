import collections 
import collections.abc
from pptx import Presentation
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Pt,Inches
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.chart.data import XySeriesData, XyChartData
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_LABEL_POSITION
import os

prs = Presentation()

blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

chart_data = ChartData()   
chart_data.categories = ['Twitter', 'Youtube', 'Instagram', 'Facebook', 'Tiktok']
chart_data.add_series('25 Jan 00:00', [451,200, 123, 345, 60])

x, y, cx, cy = Inches(1), Inches(2), Inches(5), Inches(5)

charts = slide.shapes.add_chart(XL_CHART_TYPE.PIE_EXPLODED, x, y, cx, cy, chart_data)
chart = charts.chart
# Chart Title
chart.has_title = True
chart.chart_title.has_text_frame=True
title_text_frame = chart.chart_title.text_frame
title = title_text_frame.paragraphs[0].add_run()
title.font.size = Pt(18)
title.font.bold = False
title.font.underline = True
title.text = "ini Judul"

# Chart Legend
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False
chart.legend.font.bold = True
chart.legend.font.size = Pt(10)

# Chart Plot
data_label = chart.plots[0]
data_label.has_data_labels = True
data_labels = data_label.data_labels
data_labels.position = XL_LABEL_POSITION.BEST_FIT

# Data Label In Chart Plot
data_labels.show_percentage = True
# # data_labels.show_legend_key = True
# # data_labels.show_series_name = True
data_labels.show_value = False
data_labels.show_category_name = True
# # data_labels.number_format = "dd/mm/yyyy"
# # data_labels.number_format = "Rp0.0"

# Chart Plot Styling
chart.series[0].points[1].format.fill.solid()
chart.series[0].points[1].format.fill.fore_color.rgb = RGBColor.from_string("FFC0CB")

data_labels.font.color.rgb = RGBColor.from_string("202020")
data_labels.font.size = Pt(14)

# for category in chart.plots[0].categories:
#     print(category)

# chart.plots.BarPlot.gap_width = 80
prs.save('test1.pptx')
os.startfile('test1.pptx')

label_potition_pie = ["BEST_FIT","CENTER","INSIDE_END","OUTSIDE_END"]