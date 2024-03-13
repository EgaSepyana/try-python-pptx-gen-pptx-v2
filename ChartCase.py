import collections 
import collections.abc
from pptx import Presentation
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Pt,Inches
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.chart.data import XySeriesData, XyChartData
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_LABEL_POSITION

import os

def pieChart():
    prs = Presentation()

    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    chart_data = ChartData()   
    chart_data.categories = ['Polri', 'Penembakan Ajudan Ferdy Sambo', 'Anies Baswedan', 'Daily Pilpres 2024', 'Demo Perangkat Desa']
    chart_data.add_series('25 Jan 00:00', [116.098,63.065, 17.045, 11.308,  6.628])

    x, y, cx, cy = Inches(1.2), Inches(.5), Inches(3.66), Inches(2.97)

    charts = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data)
    chart = charts.chart

    chart.has_title = False

    plotstyle = {
        "color" : ["878E91", "63B389", "F07E53","FFBC4E","C9C0AD"]
    }
    for idx, point in enumerate(chart.series[0].points):
                    col_idx = idx % len(plotstyle["color"])           
                    point.format.fill.solid()
                    point.format.fill.fore_color.rgb = RGBColor.from_string(plotstyle["color"][col_idx])

    data_label = chart.plots[0]
    data_label.has_data_labels = True
    data_labels = data_label.data_labels
    data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    data_labels.show_percentage = True
    data_labels.show_value = False
    data_labels.show_category_name = True
    data_labels.font.size = Pt(9)

    #  Chart 2

    chart_data = ChartData()   
    chart_data.categories = ['Jumlah reaksi pembicaraan vaksin secara keseluruhan', 'jumlah reaksi penolakan vaksin']
    chart_data.add_series('25 Jan 00:00', [83.797,34.854])

    x, y, cx, cy = Inches(5.08), Inches(.5), Inches(3.66), Inches(2.97)

    charts = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data)
    chart = charts.chart

    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(12)

    plotstyle = {
        "color" : ["D84E2E", "7E979B"]
    }
    for idx, point in enumerate(chart.series[0].points):
                    col_idx = idx % len(plotstyle["color"])           
                    point.format.fill.solid()
                    point.format.fill.fore_color.rgb = RGBColor.from_string(plotstyle["color"][col_idx])

    data_label = chart.plots[0]
    data_label.has_data_labels = True
    data_labels = data_label.data_labels
    data_labels.position = XL_LABEL_POSITION.INSIDE_END
    data_labels.show_percentage = True
    data_labels.show_value = False
    data_labels.show_category_name = False
    data_labels.font.size = Pt(12)
    data_labels.font.bold = True
    data_labels.font.color.rgb = RGBColor.from_string("FFFFFF")


    # Chart 3

    chart_data = ChartData()   
    chart_data.categories = ['22-30 tahun', '18-21 tahun','31-40 tahun','41-55 tahun','< 18 tahun','> 55 tahun']
    chart_data.add_series('25 Jan 00:00', [2235,1969,1852,1356,65,37])

    x, y, cx, cy = Inches(3.25), Inches(3.71), Inches(3.66), Inches(2.97)

    charts = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, chart_data)
    chart = charts.chart

    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.LEFT
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(12)

    plotstyle = {
        "color" : ["D84E2E", "CA6A28","E2772E","EF9164","F3B29A","F6CCBE"]
    }
    for idx, point in enumerate(chart.series[0].points):
                    col_idx = idx % len(plotstyle["color"])           
                    point.format.fill.solid()
                    point.format.fill.fore_color.rgb = RGBColor.from_string(plotstyle["color"][col_idx])

    data_label = chart.plots[0]
    data_label.has_data_labels = True
    data_labels = data_label.data_labels
    # data_labels.position = XL_LABEL_POSITION.CENTER
    data_labels.show_percentage = True
    # data_labels.show_value = False
    # data_labels.show_category_name = False
    data_labels.font.size = Pt(10)
    data_labels.font.bold = True
    data_labels.font.color.rgb = RGBColor.from_string("FFFFFF")


    # Chart 4
    


    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    chart_data = XyChartData()   
    categories = [1, 2, 3, 4, 5 , 6, 7, 8, 9, 10, 11, 12, 13, 14, 15]
    series1 = chart_data.add_series('FPI')
    series_data_1 = [75778 ,50969 ,4573 ,51256 ,9724 ,41547 ,24718 ,55268 ,606485 ,175234 ,75104 ,26092 ,513455 ,33622 ,12418 ]
    series2 = chart_data.add_series('NU')
    series_data_2 =  [162472 ,87099 ,38780 ,53955 ,92383 ,136720 ,62121 ,86199 ,89227 ,135517 ,157646 ,61980 ,44801 ,86957 ,37540] 
    series3 = chart_data.add_series('PKS')
    series_data_3 = [ 51572 ,38751 ,41173 ,50027 ,77237 ,18172 ,31909 ,60370 ,21602 ,61857 ,41802 ,47473 ,43129 ,50727 ,27625 ]
    series4 = chart_data.add_series('MUI')
    series_data_4 =  [80411 ,32474 ,34721 ,46544 ,39810 ,26054 ,51635 ,33774 ,30265 ,35602 ,59660 ,42404 ,48879 ,31338 ,12389 ]
    series5 = chart_data.add_series('HTI')
    series_data_5 =  [34741 ,24565 ,14032 ,67091 ,71681 ,83089 ,24735 ,25839 ,8208 ,19374 ,35087 ,22781 ,21095 ,20459 ,15676]
    series6 = chart_data.add_series('Muhammadiyah')
    series_data_6  =[ 32214 , 11293 , 18409 , 19953 , 60464 ,12198 ,24736 , 31139, 22143 , 63906 , 61484 , 37866 ,24226 , 24060 ,21555 ]
    series7 = chart_data.add_series('PAN')
    series_data_7 =  [44094 ,27760 ,18351 ,20012 ,20807 ,9113 ,8849 ,16188 ,4561 ,11279 ,11845 ,15255 ,14338 ,19948 ,6879 ]
    series8 = chart_data.add_series('PKB')
    series_data_8 = [ 12887 ,7078 ,15707 ,7346 ,34639 ,7424 ,7535 ,13368 ,4032 ,5911 ,9745 ,9170 ,17687 ,10468 ,10008 ]
    series9 = chart_data.add_series('PPP')
    series_data_9 = [ 354 ,116 ,1250 , 2714 , 8880 , 2458 , 2310 , 11122 , 3739 , 20588 , 5934 , 8774 , 10952 , 6597 , 2200 ]
    series10 = chart_data.add_series('HAMAS')
    series_data_10 = [ 371 ,195 ,141 ,315 ,616 ,1148 ,715 ,294 ,241 ,621 ,230 ,590 ,234 ,316 ,77361]


    for i in range(len(categories)):
        series1.add_data_point(x=categories[i], y=series_data_1[i])
    
    for i in range(len(categories)):
        series2.add_data_point(x=categories[i], y=series_data_2[i])
    for i in range(len(categories)):
        series3.add_data_point(x=categories[i], y=series_data_3[i])
    for i in range(len(categories)):
        series4.add_data_point(x=categories[i], y=series_data_4[i])
    
    for i in range(len(categories)):
        series5.add_data_point(x=categories[i], y=series_data_5[i])

    for i in range(len(categories)):
        series6.add_data_point(x=categories[i], y=series_data_6[i])

    for i in range(len(categories)):
        series7.add_data_point(x=categories[i], y=series_data_7[i])
    
    for i in range(len(categories)):
        series8.add_data_point(x=categories[i], y=series_data_8[i])

    for i in range(len(categories)):
        series9.add_data_point(x=categories[i], y=series_data_9[i])

    for i in range(len(categories)):
        series10.add_data_point(x=categories[i], y=series_data_10[i])
    
    # for i in [48.573 	, 38.780 ,	 41.173 	, 34.721 	, 14.032 	, 18.409 	, 18.351 	, 15.707 	, 1.250 	, 141]:
    #     series3.add_data_point(x=5, y=i)

    # for i in [ 51.256 	, 53.955 ,	 50.027 	, 46.544 	, 67.091 ,	 19.953 ,	 20.012 ,	 7.346 	, 2.714 ,	 315 ]:
    #     series4.add_data_point(x=4, y=i)

    # for i in [ 90.724 ,	 92.383 	, 77.237 	, 39.810 ,	 71.681 ,	 60.464 ,	 20.807 	, 34.639 ,	 8.880 	, 616 ]:
    #     series5.add_data_point(x=3, y=i)

    # for i in [ 41.547 	, 136.720 	, 18.172 	, 26.054 ,	 83.089 ,	 12.198 ,	 9.113 	, 7.424 	, 2.458 	, 1.148 ]:
    #     series6.add_data_point(x=2, y=i)

    # for i in [ 24.718 	, 62.121 ,	 31.909 ,	 51.635 ,	 24.735 ,	 24.736 ,	 8.849 ,	 7.535 ,	 2.310 ,	 715 ]:
    #     series7.add_data_point(x=1, y=i)

    x, y, cx, cy = Inches(0.32), Inches(0.79), Inches(9.36), Inches(5.57)

    charts = slide.shapes.add_chart(XL_CHART_TYPE.XY_SCATTER_SMOOTH_NO_MARKERS, x, y, cx, cy, chart_data)
    chart = charts.chart

    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(12)

    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    x, y, cx, cy = Inches(0.32), Inches(0.79), Inches(9.36), Inches(5.57)

    charts = slide.shapes.add_chart(XL_CHART_TYPE.XY_SCATTER_SMOOTH_NO_MARKERS, x, y, cx, cy, chart_data)
    chart = charts.chart
    
    prs.save('pieChartEx.pptx')
    os.startfile('pieChartEx.pptx')
    


if __name__ == "__main__":
    pieChart()